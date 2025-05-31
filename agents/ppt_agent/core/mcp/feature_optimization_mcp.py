"""
特性优化MCP模块

负责PPT智能体的核心功能和特性优化
"""

import os
import logging
from typing import Dict, Any, List, Optional
from pptx import Presentation
from datetime import datetime

from .base_mcp import BaseMCP

class FeatureOptimizationMCP(BaseMCP):
    """特性优化MCP，负责PPT智能体的核心功能和特性优化"""
    
    def __init__(self):
        """初始化特性优化MCP"""
        super().__init__(
            mcp_id="feature_optimization",
            name="特性优化MCP",
            description="负责PPT生成、模板应用、内容处理等核心功能"
        )
        
    def process(self, input_data: Dict[str, Any], context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        处理PPT生成相关的特性优化
        
        参数:
            input_data: 输入数据字典
            context: 上下文信息字典
            
        返回:
            处理结果字典
        """
        task_type = input_data.get("task_type")
        
        if not task_type:
            return {"status": "error", "message": "缺少task_type参数"}
            
        if task_type == "text_to_ppt":
            return self._process_text_to_ppt(input_data, context)
        elif task_type == "mindmap_to_ppt":
            return self._process_mindmap_to_ppt(input_data, context)
        elif task_type == "template_ppt":
            return self._process_template_ppt(input_data, context)
        else:
            return {"status": "error", "message": f"不支持的任务类型: {task_type}"}
    
    def _process_text_to_ppt(self, input_data: Dict[str, Any], context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        处理文本转PPT任务
        
        参数:
            input_data: 输入数据字典
            context: 上下文信息字典
            
        返回:
            处理结果字典
        """
        title = input_data.get("title", "未命名演示文稿")
        content = input_data.get("content", "")
        template_name = input_data.get("template_name", "专业简历.pptx")
        output_dir = context.get("output_dir") if context else None
        template_dir = context.get("template_dir") if context else None
        
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "output")
            os.makedirs(output_dir, exist_ok=True)
            
        if not template_dir:
            template_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "templates")
        
        # 解析内容，生成幻灯片结构
        slides_data = self._parse_content_to_slides(content)
        
        # 使用模板创建演示文稿
        template_path = os.path.join(template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            return {"status": "error", "message": f"模板文件不存在: {template_name}"}
        
        # 生成PPT
        output_path = self._generate_ppt(template_path, title, slides_data, output_dir)
        
        return {
            "status": "success",
            "data": {
                "ppt_path": output_path,
                "title": title,
                "slides_count": len(slides_data)
            }
        }
    
    def _process_mindmap_to_ppt(self, input_data: Dict[str, Any], context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        处理思维导图转PPT任务
        
        参数:
            input_data: 输入数据字典
            context: 上下文信息字典
            
        返回:
            处理结果字典
        """
        mindmap_data = input_data.get("mindmap_data", {})
        title = input_data.get("title", "思维导图演示文稿")
        template_name = input_data.get("template_name", "专业简历.pptx")
        output_dir = context.get("output_dir") if context else None
        template_dir = context.get("template_dir") if context else None
        
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "output")
            os.makedirs(output_dir, exist_ok=True)
            
        if not template_dir:
            template_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "templates")
        
        # 解析思维导图数据，生成幻灯片结构
        slides_data = self._parse_mindmap_to_slides(mindmap_data)
        
        # 使用模板创建演示文稿
        template_path = os.path.join(template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            return {"status": "error", "message": f"模板文件不存在: {template_name}"}
        
        # 生成PPT
        output_path = self._generate_ppt(template_path, title, slides_data, output_dir)
        
        return {
            "status": "success",
            "data": {
                "ppt_path": output_path,
                "title": title,
                "slides_count": len(slides_data)
            }
        }
    
    def _process_template_ppt(self, input_data: Dict[str, Any], context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        处理模板PPT任务
        
        参数:
            input_data: 输入数据字典
            context: 上下文信息字典
            
        返回:
            处理结果字典
        """
        template_name = input_data.get("template_name", "专业简历.pptx")
        title = input_data.get("title", "模板演示文稿")
        output_dir = context.get("output_dir") if context else None
        template_dir = context.get("template_dir") if context else None
        
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "output")
            os.makedirs(output_dir, exist_ok=True)
            
        if not template_dir:
            template_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "templates")
        
        # 使用模板创建演示文稿
        template_path = os.path.join(template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            return {"status": "error", "message": f"模板文件不存在: {template_name}"}
        
        try:
            # 复制模板并保存为新文件
            ppt = Presentation(template_path)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"{title}_{timestamp}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            ppt.save(output_path)
            
            return {
                "status": "success",
                "data": {
                    "ppt_path": output_path,
                    "title": title,
                    "slides_count": len(ppt.slides)
                }
            }
        except Exception as e:
            self.logger.error(f"处理模板PPT失败: {str(e)}")
            return {"status": "error", "message": f"处理模板PPT失败: {str(e)}"}
    
    def _parse_content_to_slides(self, content: str) -> List[Dict[str, Any]]:
        """
        解析文本内容，生成幻灯片结构
        
        参数:
            content: 文本内容
            
        返回:
            幻灯片数据列表
        """
        slides_data = []
        
        # 简单的解析逻辑，可以根据实际需求扩展
        lines = content.strip().split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 以#开头的行作为幻灯片标题
            if line.startswith('#'):
                if current_slide:
                    slides_data.append(current_slide)
                
                current_slide = {
                    "title": line.lstrip('#').strip(),
                    "content": []
                }
            
            # 其他行作为幻灯片内容
            elif current_slide:
                current_slide["content"].append(line)
            
            # 如果还没有幻灯片，创建第一张
            else:
                current_slide = {
                    "title": "封面",
                    "content": [line]
                }
        
        # 添加最后一张幻灯片
        if current_slide:
            slides_data.append(current_slide)
            
        # 如果没有解析出幻灯片，创建一个默认的
        if not slides_data:
            slides_data.append({
                "title": "演示文稿",
                "content": [content]
            })
        
        return slides_data
    
    def _parse_mindmap_to_slides(self, mindmap_data: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        解析思维导图数据，生成幻灯片结构
        
        参数:
            mindmap_data: 思维导图数据
            
        返回:
            幻灯片数据列表
        """
        slides_data = []
        
        # 处理根节点作为封面
        root = mindmap_data.get("root", {})
        root_title = root.get("text", "演示文稿")
        
        # 添加封面幻灯片
        slides_data.append({
            "title": root_title,
            "content": []
        })
        
        # 处理子节点作为内容幻灯片
        children = root.get("children", [])
        for child in children:
            slide_title = child.get("text", "")
            slide_content = []
            
            # 提取节点内容
            note = child.get("note", "")
            if note:
                slide_content.append(note)
                
            # 处理子节点作为幻灯片内容
            child_children = child.get("children", [])
            for sub_child in child_children:
                sub_text = sub_child.get("text", "")
                if sub_text:
                    slide_content.append(sub_text)
                    
                # 处理更深层次的内容
                sub_note = sub_child.get("note", "")
                if sub_note:
                    slide_content.append(f"  {sub_note}")
            
            # 添加内容幻灯片
            slides_data.append({
                "title": slide_title,
                "content": slide_content
            })
        
        # 添加总结幻灯片
        slides_data.append({
            "title": "总结",
            "content": ["感谢观看"]
        })
        
        return slides_data
    
    def _generate_ppt(self, template_path: str, title: str, slides_data: List[Dict[str, Any]], output_dir: str) -> str:
        """
        生成PPT文件
        
        参数:
            template_path: 模板文件路径
            title: 演示文稿标题
            slides_data: 幻灯片数据列表
            output_dir: 输出目录
            
        返回:
            生成的PPT文件路径
        """
        try:
            # 加载模板
            ppt = Presentation(template_path)
            
            # 获取模板幻灯片布局
            slide_layouts = ppt.slide_layouts
            title_slide_layout = slide_layouts[0]  # 标题幻灯片布局
            content_slide_layout = slide_layouts[1]  # 内容幻灯片布局
            
            # 保留模板幻灯片，只清除文本内容
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape.text = ""
            
            # 创建幻灯片
            for i, slide_data in enumerate(slides_data):
                slide_title = slide_data["title"]
                slide_content = slide_data["content"]
                
                # 选择布局
                layout = title_slide_layout if i == 0 else content_slide_layout
                
                # 创建幻灯片
                slide = ppt.slides.add_slide(layout)
                
                # 设置标题
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                
                # 设置内容
                if i > 0 and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_text = "\n".join(slide_content)
                    content_placeholder.text = content_text
            
            # 保存PPT
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"{title}_{timestamp}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            ppt.save(output_path)
            
            self.logger.info(f"PPT生成成功: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"PPT生成失败: {str(e)}")
            raise Exception(f"PPT生成失败: {str(e)}")
