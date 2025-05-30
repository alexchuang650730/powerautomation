"""
PPT智能体模块

负责处理PPT生成相关任务，继承自BaseAgent
集成MCP规划器和MCP头脑风暴器调用开发工具模块和已有工具
"""

import os
import logging
import uuid
from typing import Dict, Any, List, Optional
from pptx import Presentation
from datetime import datetime

from ..base.base_agent import BaseAgent
from ...agents.ppt_agent.core.mcp.mcp_planner import MCPPlanner
from ...agents.ppt_agent.core.mcp.mcp_brainstorm import MCPBrainstorm
from ...development_tools.thought_action_recorder import ThoughtActionRecorder

class PPTAgent(BaseAgent):
    """PPT智能体，负责生成和管理PPT演示文稿"""
    
    def __init__(self, agent_id: str = None):
        """
        初始化PPT智能体
        
        参数:
            agent_id: 智能体ID，如果不提供则自动生成
        """
        super().__init__(
            agent_id=agent_id,
            name="PPT智能体",
            description="省时高效的专家级PPT智能体"
        )
        self.template_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "static", "templates")
        self.output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "output")
        
        # 初始化MCP规划器和头脑风暴器
        self.mcp_planner = MCPPlanner()
        self.mcp_brainstorm = MCPBrainstorm()
        
        # 初始化思考与操作记录器
        self.recorder = ThoughtActionRecorder()
        
        # 确保输出目录存在
        os.makedirs(self.output_dir, exist_ok=True)
        
    def get_capabilities(self) -> List[str]:
        """
        获取PPT智能体能力列表
        
        返回:
            能力描述列表
        """
        return [
            "根据主题和内容生成专业PPT",
            "支持多种PPT模板",
            "自动排版和设计",
            "支持思维导图转PPT",
            "支持文本大纲转PPT"
        ]
    
    def validate_input(self, input_data: Dict[str, Any]) -> bool:
        """
        验证输入数据是否有效
        
        参数:
            input_data: 输入数据字典
            
        返回:
            数据是否有效
        """
        # 使用MCP规划器验证输入
        validation_result = self.mcp_planner.plan({
            "type": "input_validation",
            "input_data": input_data
        })
        
        if not validation_result.get("success", False):
            self.logger.error(f"输入验证失败: {validation_result.get('message', '未知错误')}")
            return False
            
        return True
    
    def process(self, input_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        处理PPT生成任务
        
        参数:
            input_data: 输入数据字典，包含任务类型和相关参数
            
        返回:
            处理结果字典
        """
        # 记录会话
        session_id = self.recorder.start_session("ppt_agent")
        self.recorder.record_thought(session_id, f"开始处理PPT生成任务: {input_data.get('task_type', '未知任务')}")
        
        # 使用MCP规划器解析任务
        self.recorder.record_thought(session_id, "使用MCP规划器解析任务")
        task_plan = self.mcp_planner.plan({
            "type": "task_parsing",
            "input_data": input_data
        })
        
        # 如果MCP规划器无法处理，尝试使用MCP头脑风暴器
        if not task_plan.get("success"):
            self.recorder.record_thought(session_id, "MCP规划器无法处理，尝试使用MCP头脑风暴器")
            task_plan = self.mcp_brainstorm.generate({
                "type": "task_parsing",
                "input_data": input_data
            })
        
        # 根据任务类型执行相应的处理
        task_type = input_data.get("task_type", "")
        
        if task_type == "text_to_ppt":
            self.recorder.record_thought(session_id, "执行文本转PPT任务")
            result = self._process_text_to_ppt(input_data, session_id)
            
        elif task_type == "mindmap_to_ppt":
            self.recorder.record_thought(session_id, "执行思维导图转PPT任务")
            result = self._process_mindmap_to_ppt(input_data, session_id)
            
        elif task_type == "template_ppt":
            self.recorder.record_thought(session_id, "执行模板PPT任务")
            result = self._process_template_ppt(input_data, session_id)
            
        else:
            self.logger.error(f"不支持的任务类型: {task_type}")
            self.recorder.record_thought(session_id, f"不支持的任务类型: {task_type}")
            result = {"error": f"不支持的任务类型: {task_type}"}
        
        # 记录操作结果
        self.recorder.record_action(session_id, "process_result", input_data, result)
        
        return result
    
    def _process_text_to_ppt(self, input_data: Dict[str, Any], session_id: str) -> Dict[str, Any]:
        """
        处理文本转PPT任务
        
        参数:
            input_data: 输入数据字典
            session_id: 会话ID
            
        返回:
            处理结果字典
        """
        # 使用MCP规划器解析文本内容
        self.recorder.record_thought(session_id, "使用MCP规划器解析文本内容")
        content_parsing = self.mcp_planner.plan({
            "type": "content_parsing",
            "input_data": input_data
        })
        
        title = input_data.get("title", "未命名演示文稿")
        content = input_data.get("content", "")
        template_name = input_data.get("template_name", "专业简历.pptx")
        
        # 解析内容，生成幻灯片结构
        self.recorder.record_thought(session_id, "解析内容，生成幻灯片结构")
        slides_data = self._parse_content_to_slides(content)
        
        # 使用模板创建演示文稿
        template_path = os.path.join(self.template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            self.recorder.record_thought(session_id, f"模板文件不存在: {template_path}")
            return {"error": f"模板文件不存在: {template_name}"}
        
        # 生成PPT
        self.recorder.record_thought(session_id, "生成PPT")
        output_path = self._generate_ppt(template_path, title, slides_data)
        
        # 记录操作结果
        result = {
            "ppt_path": output_path,
            "title": title,
            "slides_count": len(slides_data)
        }
        self.recorder.record_action(session_id, "text_to_ppt_result", input_data, result)
        
        return result
    
    def _process_mindmap_to_ppt(self, input_data: Dict[str, Any], session_id: str) -> Dict[str, Any]:
        """
        处理思维导图转PPT任务
        
        参数:
            input_data: 输入数据字典
            session_id: 会话ID
            
        返回:
            处理结果字典
        """
        # 使用MCP规划器解析思维导图数据
        self.recorder.record_thought(session_id, "使用MCP规划器解析思维导图数据")
        mindmap_parsing = self.mcp_planner.plan({
            "type": "mindmap_parsing",
            "input_data": input_data
        })
        
        mindmap_data = input_data["mindmap_data"]
        title = input_data.get("title", "思维导图演示文稿")
        template_name = input_data.get("template_name", "专业简历.pptx")
        
        # 解析思维导图数据，生成幻灯片结构
        self.recorder.record_thought(session_id, "解析思维导图数据，生成幻灯片结构")
        slides_data = self._parse_mindmap_to_slides(mindmap_data)
        
        # 使用模板创建演示文稿
        template_path = os.path.join(self.template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            self.recorder.record_thought(session_id, f"模板文件不存在: {template_path}")
            return {"error": f"模板文件不存在: {template_name}"}
        
        # 生成PPT
        self.recorder.record_thought(session_id, "生成PPT")
        output_path = self._generate_ppt(template_path, title, slides_data)
        
        # 记录操作结果
        result = {
            "ppt_path": output_path,
            "title": title,
            "slides_count": len(slides_data)
        }
        self.recorder.record_action(session_id, "mindmap_to_ppt_result", input_data, result)
        
        return result
    
    def _process_template_ppt(self, input_data: Dict[str, Any], session_id: str) -> Dict[str, Any]:
        """
        处理模板PPT任务
        
        参数:
            input_data: 输入数据字典
            session_id: 会话ID
            
        返回:
            处理结果字典
        """
        template_name = input_data["template_name"]
        title = input_data.get("title", "模板演示文稿")
        
        # 使用模板创建演示文稿
        template_path = os.path.join(self.template_dir, template_name)
        if not os.path.exists(template_path):
            self.logger.error(f"模板文件不存在: {template_path}")
            self.recorder.record_thought(session_id, f"模板文件不存在: {template_path}")
            return {"error": f"模板文件不存在: {template_name}"}
        
        # 复制模板并保存为新文件
        self.recorder.record_thought(session_id, "复制模板并保存为新文件")
        ppt = Presentation(template_path)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"{title}_{timestamp}.pptx"
        output_path = os.path.join(self.output_dir, output_filename)
        ppt.save(output_path)
        
        # 记录操作结果
        result = {
            "ppt_path": output_path,
            "title": title,
            "slides_count": len(ppt.slides)
        }
        self.recorder.record_action(session_id, "template_ppt_result", input_data, result)
        
        return result
    
    def _parse_content_to_slides(self, content: str) -> List[Dict[str, Any]]:
        """
        解析文本内容，生成幻灯片结构
        
        参数:
            content: 文本内容
            
        返回:
            幻灯片数据列表
        """
        slides_data = []
        
        # 简单的解析逻辑，可以根据实际需求扩展
        lines = content.strip().split('\n')
        current_slide = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 以#开头的行作为幻灯片标题
            if line.startswith('#'):
                if current_slide:
                    slides_data.append(current_slide)
                
                current_slide = {
                    "title": line.lstrip('#').strip(),
                    "content": []
                }
            
            # 其他行作为幻灯片内容
            elif current_slide:
                current_slide["content"].append(line)
            
            # 如果还没有幻灯片，创建第一张
            else:
                current_slide = {
                    "title": "封面",
                    "content": [line]
                }
        
        # 添加最后一张幻灯片
        if current_slide:
            slides_data.append(current_slide)
            
        # 如果没有解析出幻灯片，创建一个默认的
        if not slides_data:
            slides_data.append({
                "title": "演示文稿",
                "content": [content]
            })
        
        return slides_data
    
    def _parse_mindmap_to_slides(self, mindmap_data: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        解析思维导图数据，生成幻灯片结构
        
        参数:
            mindmap_data: 思维导图数据
            
        返回:
            幻灯片数据列表
        """
        slides_data = []
        
        # 处理根节点作为封面
        root = mindmap_data.get("root", {})
        root_title = root.get("text", "演示文稿")
        
        # 添加封面幻灯片
        slides_data.append({
            "title": root_title,
            "content": []
        })
        
        # 处理子节点作为内容幻灯片
        children = root.get("children", [])
        for child in children:
            slide_title = child.get("text", "")
            slide_content = []
            
            # 提取节点内容
            note = child.get("note", "")
            if note:
                slide_content.append(note)
                
            # 处理子节点作为幻灯片内容
            child_children = child.get("children", [])
            for sub_child in child_children:
                sub_text = sub_child.get("text", "")
                if sub_text:
                    slide_content.append(sub_text)
                    
                # 处理更深层次的内容
                sub_note = sub_child.get("note", "")
                if sub_note:
                    slide_content.append(f"  {sub_note}")
            
            # 添加内容幻灯片
            slides_data.append({
                "title": slide_title,
                "content": slide_content
            })
        
        # 添加总结幻灯片
        slides_data.append({
            "title": "总结",
            "content": ["感谢观看"]
        })
        
        return slides_data
    
    def _generate_ppt(self, template_path: str, title: str, slides_data: List[Dict[str, Any]]) -> str:
        """
        生成PPT文件
        
        参数:
            template_path: 模板文件路径
            title: 演示文稿标题
            slides_data: 幻灯片数据列表
            
        返回:
            生成的PPT文件路径
        """
        try:
            # 加载模板
            ppt = Presentation(template_path)
            
            # 获取模板幻灯片布局
            slide_layouts = ppt.slide_layouts
            title_slide_layout = slide_layouts[0]  # 标题幻灯片布局
            content_slide_layout = slide_layouts[1]  # 内容幻灯片布局
            
            # 清除模板中的示例幻灯片
            for _ in range(len(ppt.slides)):
                rId = ppt.slides._sldIdLst[0].rId
                ppt.part.drop_rel(rId)
                del ppt.slides._sldIdLst[0]
            
            # 创建幻灯片
            for i, slide_data in enumerate(slides_data):
                slide_title = slide_data["title"]
                slide_content = slide_data["content"]
                
                # 选择布局
                layout = title_slide_layout if i == 0 else content_slide_layout
                
                # 创建幻灯片
                slide = ppt.slides.add_slide(layout)
                
                # 设置标题
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                
                # 设置内容
                if i > 0 and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_text = "\n".join(slide_content)
                    content_placeholder.text = content_text
            
            # 保存PPT
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"{title}_{timestamp}.pptx"
            output_path = os.path.join(self.output_dir, output_filename)
            ppt.save(output_path)
            
            self.logger.info(f"PPT生成成功: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"PPT生成失败: {str(e)}")
            raise Exception(f"PPT生成失败: {str(e)}")
